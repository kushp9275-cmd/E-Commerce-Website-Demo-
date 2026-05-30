"""
Generate a professional PowerPoint presentation for the Mart E-Commerce project.
Uses python-pptx with custom shapes, gradients, and diagrams.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ─────────────────────────────────────────────────────────────────────────────
# COLOR PALETTE
# ─────────────────────────────────────────────────────────────────────────────
BG_DARK        = RGBColor(0x0F, 0x17, 0x2A)   # Deep navy
BG_CARD        = RGBColor(0x1E, 0x29, 0x3B)   # Card dark
ACCENT_PRIMARY = RGBColor(0x63, 0x66, 0xF1)   # Indigo-500
ACCENT_CYAN    = RGBColor(0x22, 0xD3, 0xEE)   # Cyan-400
ACCENT_EMERALD = RGBColor(0x10, 0xB9, 0x81)   # Emerald-500
ACCENT_AMBER   = RGBColor(0xF5, 0x9E, 0x0B)   # Amber-500
ACCENT_ROSE    = RGBColor(0xF4, 0x3F, 0x5E)   # Rose-500
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY     = RGBColor(0xCB, 0xD5, 0xE1)   # Slate-300
MID_GRAY       = RGBColor(0x94, 0xA3, 0xB8)   # Slate-400
DARK_GRAY      = RGBColor(0x47, 0x55, 0x69)   # Slate-600
GRADIENT_START = RGBColor(0x63, 0x66, 0xF1)
GRADIENT_END   = RGBColor(0x8B, 0x5C, 0xF6)   # Violet-500

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ─────────────────────────────────────────────────────────────────────────────
# HELPER FUNCTIONS
# ─────────────────────────────────────────────────────────────────────────────

def set_slide_bg(slide, color):
    """Set solid background color on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    """Add a shape with optional fill and border."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.background()  # transparent by default
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def set_text(shape, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    """Set text on a shape's text frame."""
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    """Add a text box to the slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multi_text_box(slide, left, top, width, height, lines, font_name='Calibri'):
    """Add a text box with multiple styled lines. Each line is a dict with text, size, color, bold, alignment."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line.get('text', '')
        p.font.size = Pt(line.get('size', 18))
        p.font.color.rgb = line.get('color', WHITE)
        p.font.bold = line.get('bold', False)
        p.font.name = font_name
        p.alignment = line.get('alignment', PP_ALIGN.LEFT)
        p.space_after = Pt(line.get('space_after', 6))
        p.space_before = Pt(line.get('space_before', 0))
    return txBox

def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None, corner_radius=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_circle(slide, left, top, size, fill_color, line_color=None):
    """Add a circle (oval)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    return shape

def add_arrow_connector(slide, x1, y1, x2, y2, color=ACCENT_CYAN, width=Pt(2)):
    """Add a straight connector line between two points."""
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1 = straight
    connector.line.color.rgb = color
    connector.line.width = width
    return connector

def add_accent_bar(slide, left, top, width, height, color=ACCENT_PRIMARY):
    """Add a small accent bar/line decoration."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar

def add_section_number(slide, left, top, number, color=ACCENT_PRIMARY):
    """Add a circled section number."""
    circle = add_circle(slide, left, top, Inches(0.6), color)
    set_text(circle, str(number), font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    circle.text_frame.word_wrap = False
    return circle

def add_decorative_dots(slide, left, top, rows, cols, spacing, color, size=Inches(0.06)):
    """Add a grid of decorative dots."""
    for r in range(rows):
        for c in range(cols):
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                left + c * spacing,
                top + r * spacing,
                size, size
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = color
            dot.line.fill.background()

def add_gradient_bar(slide, left, top, width, height):
    """Add a gradient accent bar at top/bottom of slide."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_PRIMARY
    bar.line.fill.background()
    return bar


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE BUILDERS
# ─────────────────────────────────────────────────────────────────────────────

def build_title_slide(prs):
    """Slide 1 — Title Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_DARK)

    # Top gradient bar
    add_gradient_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06))

    # Decorative dots top-right
    add_decorative_dots(slide, Inches(10.5), Inches(0.8), 5, 5, Inches(0.25), RGBColor(0x63, 0x66, 0xF1))

    # Decorative dots bottom-left
    add_decorative_dots(slide, Inches(0.5), Inches(5.5), 4, 4, Inches(0.25), RGBColor(0x22, 0xD3, 0xEE))

    # Large decorative circle (background)
    big_circle = add_circle(slide, Inches(8.5), Inches(-1.5), Inches(10), RGBColor(0x1A, 0x22, 0x38))
    big_circle.fill.fore_color.rgb = RGBColor(0x1A, 0x22, 0x38)

    # Accent line
    add_accent_bar(slide, Inches(1.5), Inches(2.3), Inches(1.2), Inches(0.06), ACCENT_PRIMARY)

    # Title
    add_text_box(slide, Inches(1.5), Inches(2.5), Inches(8), Inches(1.2),
                 "MART — E-Commerce Platform", font_size=42, color=WHITE, bold=True, font_name='Calibri')

    # Subtitle
    add_multi_text_box(slide, Inches(1.5), Inches(3.7), Inches(7), Inches(1.5), [
        {'text': 'A Full-Stack Digital Storefront', 'size': 24, 'color': ACCENT_CYAN, 'bold': False, 'space_after': 8},
        {'text': 'Modern glassmorphic design • Razorpay payments • Real-time inventory', 'size': 14, 'color': MID_GRAY, 'bold': False, 'space_after': 4},
        {'text': 'Admin dashboard • Order tracking • Role-based access', 'size': 14, 'color': MID_GRAY, 'bold': False},
    ])

    # Divider line
    add_accent_bar(slide, Inches(1.5), Inches(5.5), Inches(3), Inches(0.02), DARK_GRAY)

    # Author / Date
    add_multi_text_box(slide, Inches(1.5), Inches(5.7), Inches(5), Inches(1), [
        {'text': 'Presented by: Kush Patel', 'size': 16, 'color': LIGHT_GRAY, 'bold': True, 'space_after': 4},
        {'text': 'Full-Stack Web Development Project  •  2026', 'size': 12, 'color': MID_GRAY, 'bold': False},
    ])

    # Bottom gradient bar
    add_gradient_bar(slide, Inches(0), Inches(7.44), SLIDE_WIDTH, Inches(0.06))


def build_agenda_slide(prs):
    """Slide 2 — Agenda / Table of Contents"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_gradient_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.04))

    # Section header
    add_accent_bar(slide, Inches(0.8), Inches(0.5), Inches(0.8), Inches(0.05), ACCENT_PRIMARY)
    add_text_box(slide, Inches(0.8), Inches(0.65), Inches(4), Inches(0.6),
                 "AGENDA", font_size=28, color=WHITE, bold=True)

    agenda_items = [
        ("01", "Project Overview", "What Mart is and the problem it solves", ACCENT_PRIMARY),
        ("02", "Technology Stack", "Backend, frontend, database & integrations", ACCENT_CYAN),
        ("03", "System Architecture", "Data flow from user to database", ACCENT_EMERALD),
        ("04", "Entity-Relationship Diagram", "Database schema and table relationships", ACCENT_AMBER),
        ("05", "Key Features & USPs", "What sets Mart apart from competitors", ACCENT_ROSE),
        ("06", "Code Highlights", "Key implementation patterns", ACCENT_PRIMARY),
        ("07", "Thank You", "Questions & discussion", ACCENT_CYAN),
    ]

    start_y = Inches(1.6)
    for i, (num, title, desc, color) in enumerate(agenda_items):
        y = start_y + Inches(i * 0.75)
        
        # Number circle
        circle = add_circle(slide, Inches(1.0), y + Inches(0.05), Inches(0.5), color)
        set_text(circle, num, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        
        # Title
        add_text_box(slide, Inches(1.8), y, Inches(4), Inches(0.4),
                     title, font_size=18, color=WHITE, bold=True)
        # Description
        add_text_box(slide, Inches(1.8), y + Inches(0.35), Inches(5), Inches(0.3),
                     desc, font_size=12, color=MID_GRAY)

    # Decorative dots
    add_decorative_dots(slide, Inches(9), Inches(2), 6, 6, Inches(0.3), RGBColor(0x2A, 0x35, 0x50))

    # Right side — vertical line
    add_accent_bar(slide, Inches(7.8), Inches(1.6), Inches(0.03), Inches(5.2), DARK_GRAY)

    # Right side summary text
    add_multi_text_box(slide, Inches(8.2), Inches(2.5), Inches(4.5), Inches(3), [
        {'text': '7 Sections', 'size': 40, 'color': ACCENT_PRIMARY, 'bold': True, 'space_after': 12},
        {'text': 'Covering the complete', 'size': 16, 'color': LIGHT_GRAY, 'bold': False, 'space_after': 2},
        {'text': 'architecture, design,', 'size': 16, 'color': LIGHT_GRAY, 'bold': False, 'space_after': 2},
        {'text': 'and unique value of', 'size': 16, 'color': LIGHT_GRAY, 'bold': False, 'space_after': 2},
        {'text': 'the Mart platform.', 'size': 16, 'color': LIGHT_GRAY, 'bold': False},
    ])

    add_gradient_bar(slide, Inches(0), Inches(7.46), SLIDE_WIDTH, Inches(0.04))


def build_overview_slide(prs):
    """Slide 3 — Project Overview"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_gradient_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.04))

    # Section number & title
    add_section_number(slide, Inches(0.8), Inches(0.4), "01", ACCENT_PRIMARY)
    add_text_box(slide, Inches(1.6), Inches(0.4), Inches(5), Inches(0.6),
                 "PROJECT OVERVIEW", font_size=28, color=WHITE, bold=True)
    add_accent_bar(slide, Inches(1.6), Inches(1.05), Inches(1.5), Inches(0.04), ACCENT_PRIMARY)

    # Left column — What is Mart?
    card1 = add_rounded_rect(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(2.5), BG_CARD, DARK_GRAY)
    add_multi_text_box(slide, Inches(0.9), Inches(1.7), Inches(5.3), Inches(2.2), [
        {'text': '🛒  What is Mart?', 'size': 20, 'color': ACCENT_CYAN, 'bold': True, 'space_after': 12},
        {'text': 'Mart is a full-stack e-commerce platform built with Flask and SQLite,', 'size': 13, 'color': LIGHT_GRAY, 'space_after': 4},
        {'text': 'featuring a premium glassmorphic UI, integrated Razorpay payments,', 'size': 13, 'color': LIGHT_GRAY, 'space_after': 4},
        {'text': 'real-time inventory management, and a powerful admin dashboard.', 'size': 13, 'color': LIGHT_GRAY, 'space_after': 10},
        {'text': 'It provides a complete shopping experience — from browsing and', 'size': 13, 'color': LIGHT_GRAY, 'space_after': 4},
        {'text': 'cart management to checkout and order tracking.', 'size': 13, 'color': LIGHT_GRAY},
    ])

    # Right column — Key Metrics
    metrics = [
        ("996", "Lines of Python", ACCENT_PRIMARY),
        ("11", "HTML Templates", ACCENT_CYAN),
        ("5", "DB Tables", ACCENT_EMERALD),
        ("15+", "API Routes", ACCENT_AMBER),
    ]

    start_x = Inches(7)
    for i, (value, label, color) in enumerate(metrics):
        row = i // 2
        col = i % 2
        x = start_x + col * Inches(2.8)
        y = Inches(1.5) + row * Inches(1.4)
        card = add_rounded_rect(slide, x, y, Inches(2.5), Inches(1.2), BG_CARD, color)
        add_multi_text_box(slide, x + Inches(0.25), y + Inches(0.15), Inches(2), Inches(0.9), [
            {'text': value, 'size': 32, 'color': color, 'bold': True, 'alignment': PP_ALIGN.CENTER, 'space_after': 2},
            {'text': label, 'size': 12, 'color': LIGHT_GRAY, 'alignment': PP_ALIGN.CENTER},
        ])

    # Bottom section — Problem & Solution
    card2 = add_rounded_rect(slide, Inches(0.6), Inches(4.3), Inches(5.8), Inches(2.8), BG_CARD, DARK_GRAY)
    add_multi_text_box(slide, Inches(0.9), Inches(4.5), Inches(5.3), Inches(2.5), [
        {'text': '🎯  Problem Statement', 'size': 18, 'color': ACCENT_ROSE, 'bold': True, 'space_after': 8},
        {'text': '• Most student e-commerce projects lack real payment integration', 'size': 12, 'color': LIGHT_GRAY, 'space_after': 4},
        {'text': '• No admin panel for real-time inventory/order management', 'size': 12, 'color': LIGHT_GRAY, 'space_after': 4},
        {'text': '• Poor UI/UX with no responsive or modern design patterns', 'size': 12, 'color': LIGHT_GRAY, 'space_after': 4},
        {'text': '• No order tracking or delivery status updates', 'size': 12, 'color': LIGHT_GRAY, 'space_after': 10},
        {'text': '✅  Mart solves all of these with a production-grade approach.', 'size': 13, 'color': ACCENT_EMERALD, 'bold': True},
    ])

    card3 = add_rounded_rect(slide, Inches(7), Inches(4.3), Inches(5.8), Inches(2.8), BG_CARD, DARK_GRAY)
    add_multi_text_box(slide, Inches(7.3), Inches(4.5), Inches(5.3), Inches(2.5), [
        {'text': '🚀  Project Goals', 'size': 18, 'color': ACCENT_EMERALD, 'bold': True, 'space_after': 8},
        {'text': '• Deliver a production-ready e-commerce experience', 'size': 12, 'color': LIGHT_GRAY, 'space_after': 4},
        {'text': '• Integrate real payment gateway (Razorpay)', 'size': 12, 'color': LIGHT_GRAY, 'space_after': 4},
        {'text': '• Build role-based access (User / Admin)', 'size': 12, 'color': LIGHT_GRAY, 'space_after': 4},
        {'text': '• Implement end-to-end order lifecycle tracking', 'size': 12, 'color': LIGHT_GRAY, 'space_after': 4},
        {'text': '• Create a beautiful, responsive glassmorphic UI', 'size': 12, 'color': LIGHT_GRAY},
    ])

    add_gradient_bar(slide, Inches(0), Inches(7.46), SLIDE_WIDTH, Inches(0.04))


def build_tech_stack_slide(prs):
    """Slide 4 — Technology Stack"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_gradient_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.04))

    add_section_number(slide, Inches(0.8), Inches(0.4), "02", ACCENT_CYAN)
    add_text_box(slide, Inches(1.6), Inches(0.4), Inches(5), Inches(0.6),
                 "TECHNOLOGY STACK", font_size=28, color=WHITE, bold=True)
    add_accent_bar(slide, Inches(1.6), Inches(1.05), Inches(1.5), Inches(0.04), ACCENT_CYAN)

    tech_cards = [
        ("🐍", "BACKEND", "Flask (Python 3)", [
            "Lightweight WSGI framework",
            "Jinja2 template engine",
            "Werkzeug password hashing",
            "Session-based authentication",
            "Gunicorn for production",
        ], ACCENT_PRIMARY),
        ("🎨", "FRONTEND", "HTML5 / CSS3 / JS", [
            "Vanilla CSS — no frameworks",
            "Glassmorphic design system",
            "Inter font from Google Fonts",
            "Dark/Light theme toggle",
            "Responsive mobile layout",
        ], ACCENT_CYAN),
        ("🗄️", "DATABASE", "SQLite", [
            "Zero-configuration setup",
            "5 normalized tables",
            "Auto-migration on startup",
            "Row-factory dict access",
            "Portable single-file DB",
        ], ACCENT_EMERALD),
        ("💳", "PAYMENTS", "Razorpay SDK", [
            "Razorpay payment gateway",
            "Signature verification",
            "Cash-on-Delivery fallback",
            "Secure .env key storage",
            "Graceful error handling",
        ], ACCENT_AMBER),
        ("🔧", "DEV TOOLS", "Environment", [
            "python-dotenv for secrets",
            "Git version control",
            "Render cloud deployment",
            "Modular file structure",
            "Auto DB seed on startup",
        ], ACCENT_ROSE),
    ]

    start_x = Inches(0.4)
    card_w = Inches(2.35)
    gap = Inches(0.17)

    for i, (emoji, category, title, bullets, color) in enumerate(tech_cards):
        x = start_x + i * (card_w + gap)
        y = Inches(1.5)

        # Card background
        card = add_rounded_rect(slide, x, y, card_w, Inches(5.5), BG_CARD, color)

        # Emoji icon circle
        icon_circle = add_circle(slide, x + Inches(0.75), y + Inches(0.25), Inches(0.7), color)
        set_text(icon_circle, emoji, font_size=24, color=WHITE, alignment=PP_ALIGN.CENTER)

        # Category label
        add_text_box(slide, x + Inches(0.1), y + Inches(1.1), card_w - Inches(0.2), Inches(0.35),
                     category, font_size=11, color=color, bold=True, alignment=PP_ALIGN.CENTER)

        # Title
        add_text_box(slide, x + Inches(0.1), y + Inches(1.4), card_w - Inches(0.2), Inches(0.4),
                     title, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        # Separator
        add_accent_bar(slide, x + Inches(0.5), y + Inches(1.9), card_w - Inches(1), Inches(0.02), color)

        # Bullet points
        lines = []
        for b in bullets:
            lines.append({'text': f'▸  {b}', 'size': 10, 'color': LIGHT_GRAY, 'space_after': 6})
        add_multi_text_box(slide, x + Inches(0.15), y + Inches(2.1), card_w - Inches(0.3), Inches(3), lines)

    add_gradient_bar(slide, Inches(0), Inches(7.46), SLIDE_WIDTH, Inches(0.04))


def build_architecture_slide(prs):
    """Slide 5 — System Architecture / Data Flow Diagram"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_gradient_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.04))

    add_section_number(slide, Inches(0.8), Inches(0.4), "03", ACCENT_EMERALD)
    add_text_box(slide, Inches(1.6), Inches(0.4), Inches(6), Inches(0.6),
                 "SYSTEM ARCHITECTURE & DATA FLOW", font_size=28, color=WHITE, bold=True)
    add_accent_bar(slide, Inches(1.6), Inches(1.05), Inches(2), Inches(0.04), ACCENT_EMERALD)

    # ── Define the architecture boxes ──
    # Row 1: User / Browser
    layers = [
        # (label, x, y, w, h, color, sub_items)
        ("👤  USER / BROWSER", Inches(0.5), Inches(1.6), Inches(12.3), Inches(0.9), ACCENT_PRIMARY,
         "HTML5 Forms  •  JavaScript  •  CSS3 Glassmorphic UI  •  Theme Toggle  •  Responsive Layout"),

        ("🌐  FLASK WEB SERVER  (app.py)", Inches(0.5), Inches(3.0), Inches(12.3), Inches(0.9), ACCENT_CYAN,
         "Route Handlers  •  Jinja2 Rendering  •  Session Management  •  Input Validation  •  Error Handling"),

        ("⚙️  BUSINESS LOGIC LAYER", Inches(0.5), Inches(4.4), Inches(12.3), Inches(0.9), ACCENT_EMERALD,
         "Auth (Login/Register)  •  Cart CRUD  •  Checkout  •  Order Lifecycle  •  Admin CRUD  •  Stock Mgmt"),

        ("🗄️  DATA LAYER  (SQLite — mart.db)", Inches(0.5), Inches(5.8), Inches(8), Inches(0.9), ACCENT_AMBER,
         "users  •  items  •  cart  •  orders  •  order_items"),
    ]

    for label, x, y, w, h, color, sub in layers:
        card = add_rounded_rect(slide, x, y, w, h, BG_CARD, color)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.08), w - Inches(0.4), Inches(0.4),
                     label, font_size=15, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.45), w - Inches(0.4), Inches(0.4),
                     sub, font_size=10, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    # Razorpay external box
    rp_card = add_rounded_rect(slide, Inches(9), Inches(5.8), Inches(3.8), Inches(0.9), BG_CARD, ACCENT_ROSE)
    add_text_box(slide, Inches(9.2), Inches(5.88), Inches(3.4), Inches(0.4),
                 "💳  RAZORPAY API", font_size=15, color=ACCENT_ROSE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(9.2), Inches(6.25), Inches(3.4), Inches(0.4),
                 "Order Create  •  Payment Verify  •  Signature Check", font_size=10, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    # Arrows (vertical connectors between layers)
    arrow_x_positions = [Inches(6.6)]
    arrow_pairs = [
        (Inches(2.5), Inches(3.0)),  # User -> Flask
        (Inches(3.9), Inches(4.4)),  # Flask -> Logic
        (Inches(5.3), Inches(5.8)),  # Logic -> Data
    ]

    for ax in arrow_x_positions:
        for y_start, y_end in arrow_pairs:
            # Down arrow
            arrow_down = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, ax - Inches(0.15), y_start, Inches(0.3), y_end - y_start)
            arrow_down.fill.solid()
            arrow_down.fill.fore_color.rgb = ACCENT_CYAN
            arrow_down.line.fill.background()
            # Up arrow (response)
            arrow_up = slide.shapes.add_shape(MSO_SHAPE.UP_ARROW, ax + Inches(0.35), y_start, Inches(0.3), y_end - y_start)
            arrow_up.fill.solid()
            arrow_up.fill.fore_color.rgb = ACCENT_EMERALD
            arrow_up.line.fill.background()

    # Labels for arrows
    add_text_box(slide, Inches(5.1), Inches(2.55), Inches(1.2), Inches(0.3),
                 "Request ▼", font_size=9, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(7.2), Inches(2.55), Inches(1.2), Inches(0.3),
                 "▲ Response", font_size=9, color=ACCENT_EMERALD, alignment=PP_ALIGN.CENTER)

    # Legend
    add_multi_text_box(slide, Inches(0.5), Inches(7.0), Inches(12), Inches(0.4), [
        {'text': '▸ Request Flow (Blue ▼)     ▸ Response Flow (Green ▲)     ▸ External API calls to Razorpay for online payments', 
         'size': 10, 'color': MID_GRAY, 'alignment': PP_ALIGN.LEFT},
    ])

    add_gradient_bar(slide, Inches(0), Inches(7.46), SLIDE_WIDTH, Inches(0.04))


def build_er_diagram_slide(prs):
    """Slide 6 — Entity-Relationship Diagram"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_gradient_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.04))

    add_section_number(slide, Inches(0.8), Inches(0.4), "04", ACCENT_AMBER)
    add_text_box(slide, Inches(1.6), Inches(0.4), Inches(7), Inches(0.6),
                 "ENTITY-RELATIONSHIP DIAGRAM", font_size=28, color=WHITE, bold=True)
    add_accent_bar(slide, Inches(1.6), Inches(1.05), Inches(2), Inches(0.04), ACCENT_AMBER)

    # Define tables
    tables = [
        ("USERS", [
            "🔑 id  INTEGER  PK",
            "username  VARCHAR(50)",
            "email  VARCHAR(100)  UNIQUE",
            "password_hash  VARCHAR(255)",
            "mobile_no  VARCHAR(20)",
            "address  TEXT",
            "profile_pic  VARCHAR(255)",
            "role  VARCHAR(20)",
            "created_at  TIMESTAMP",
        ], Inches(0.3), Inches(1.5), Inches(2.8), ACCENT_PRIMARY),

        ("ITEMS", [
            "🔑 id  INTEGER  PK",
            "name  VARCHAR(100)",
            "description  TEXT",
            "price  DECIMAL(10,2)",
            "image_url  VARCHAR(255)",
            "category  VARCHAR(50)",
            "stock  INT",
        ], Inches(5.2), Inches(1.5), Inches(2.8), ACCENT_CYAN),

        ("CART", [
            "🔑 id  INTEGER  PK",
            "user_id  INT  →  FK",
            "item_id  INT  →  FK",
            "quantity  INT",
        ], Inches(10.2), Inches(1.5), Inches(2.8), ACCENT_EMERALD),

        ("ORDERS", [
            "🔑 id  INTEGER  PK",
            "user_id  INT  →  FK",
            "total_price  DECIMAL(10,2)",
            "razorpay_order_id  VARCHAR",
            "razorpay_payment_id  VARCHAR",
            "payment_method  VARCHAR(20)",
            "status  VARCHAR(30)",
            "delivery_address  TEXT",
            "order_date  TIMESTAMP",
        ], Inches(0.3), Inches(4.6), Inches(3.5), ACCENT_AMBER),

        ("ORDER_ITEMS", [
            "🔑 id  INTEGER  PK",
            "order_id  INT  →  FK",
            "item_id  INT  →  FK",
            "item_name  VARCHAR(100)",
            "quantity  INT",
            "price_at_time  DECIMAL(10,2)",
            "image_url  VARCHAR(255)",
        ], Inches(5.2), Inches(4.6), Inches(3.5), ACCENT_ROSE),
    ]

    for tbl_name, columns, x, y, w, color in tables:
        row_h = Inches(0.23)
        total_h = Inches(0.45) + len(columns) * row_h + Inches(0.15)

        card = add_rounded_rect(slide, x, y, w, total_h, BG_CARD, color)

        # Header bar
        header = add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.4), fill_color=color)
        header.line.fill.background()
        set_text(header, f"  {tbl_name}", font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)

        # Columns
        lines = [{'text': col, 'size': 9, 'color': LIGHT_GRAY, 'space_after': 2} for col in columns]
        add_multi_text_box(slide, x + Inches(0.1), y + Inches(0.45), w - Inches(0.2), total_h - Inches(0.5), lines)

    # Relationship labels
    relationships = [
        ("users.id ──── 1:N ────▶ cart.user_id", Inches(9.5), Inches(4.1)),
        ("users.id ──── 1:N ────▶ orders.user_id", Inches(0.3), Inches(4.25)),
        ("items.id ──── 1:N ────▶ cart.item_id", Inches(9.5), Inches(4.35)),
        ("orders.id ── 1:N ──▶ order_items.order_id", Inches(5.2), Inches(4.25)),
    ]

    for text, x, y in relationships:
        add_text_box(slide, x, y, Inches(3.5), Inches(0.25),
                     text, font_size=9, color=ACCENT_CYAN, bold=False, alignment=PP_ALIGN.LEFT)

    add_gradient_bar(slide, Inches(0), Inches(7.46), SLIDE_WIDTH, Inches(0.04))


def build_usp_slide(prs):
    """Slide 7 — Key Features & USPs"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_gradient_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.04))

    add_section_number(slide, Inches(0.8), Inches(0.4), "05", ACCENT_ROSE)
    add_text_box(slide, Inches(1.6), Inches(0.4), Inches(8), Inches(0.6),
                 "WHAT SETS MART APART?", font_size=28, color=WHITE, bold=True)
    add_accent_bar(slide, Inches(1.6), Inches(1.05), Inches(2), Inches(0.04), ACCENT_ROSE)

    add_text_box(slide, Inches(1.6), Inches(1.15), Inches(10), Inches(0.4),
                 "Features that most student / demo e-commerce projects do NOT have", font_size=13, color=MID_GRAY)

    usps = [
        ("💳", "Real Payment Gateway", "Integrated Razorpay with signature verification — not a mock or placeholder. Supports both online pay and COD.", ACCENT_PRIMARY),
        ("📦", "Live Order Tracking", "Multi-stage order lifecycle: Confirmed → Shipped → Out for Delivery → Delivered, with a visual tracker.", ACCENT_CYAN),
        ("👨‍💼", "Full Admin Dashboard", "Admins can add/edit/delete products, manage stock, view revenue stats, and update order statuses in real-time.", ACCENT_EMERALD),
        ("🎨", "Glassmorphic UI", "Premium dark-mode design with glassmorphism, micro-animations, Inter typography — not Bootstrap or plain HTML.", ACCENT_AMBER),
        ("🛡️", "Role-Based Access", "Separate User and Admin flows with session-based auth, password hashing (Werkzeug), and login role enforcement.", ACCENT_ROSE),
        ("📊", "Smart Inventory", "Real-time stock tracking, low-stock alerts on admin dashboard, auto-decrement on purchase, and out-of-stock badges.", RGBColor(0xA7, 0x8B, 0xFA)),
    ]

    cols = 3
    rows = 2
    card_w = Inches(3.9)
    card_h = Inches(2.5)
    gap_x = Inches(0.25)
    gap_y = Inches(0.25)
    start_x = Inches(0.45)
    start_y = Inches(1.7)

    for i, (emoji, title, desc, color) in enumerate(usps):
        row = i // cols
        col = i % cols
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        card = add_rounded_rect(slide, x, y, card_w, card_h, BG_CARD, color)

        # Icon
        icon_bg = add_rounded_rect(slide, x + Inches(0.2), y + Inches(0.2), Inches(0.55), Inches(0.55), color)
        set_text(icon_bg, emoji, font_size=20, color=WHITE, alignment=PP_ALIGN.CENTER)

        # Title
        add_text_box(slide, x + Inches(0.9), y + Inches(0.2), card_w - Inches(1.1), Inches(0.45),
                     title, font_size=16, color=WHITE, bold=True)

        # Separator
        add_accent_bar(slide, x + Inches(0.2), y + Inches(0.9), card_w - Inches(0.4), Inches(0.02), color)

        # Description
        add_text_box(slide, x + Inches(0.2), y + Inches(1.05), card_w - Inches(0.4), Inches(1.3),
                     desc, font_size=11, color=LIGHT_GRAY)

    add_gradient_bar(slide, Inches(0), Inches(7.46), SLIDE_WIDTH, Inches(0.04))


def build_comparison_slide(prs):
    """Slide 8 — Mart vs Other E-Commerce Projects (comparison table)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_gradient_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.04))

    add_section_number(slide, Inches(0.8), Inches(0.4), "05", ACCENT_ROSE)
    add_text_box(slide, Inches(1.6), Inches(0.4), Inches(8), Inches(0.6),
                 "MART  vs  TYPICAL E-COMMERCE PROJECTS", font_size=26, color=WHITE, bold=True)
    add_accent_bar(slide, Inches(1.6), Inches(1.05), Inches(2.5), Inches(0.04), ACCENT_ROSE)

    # Table headers
    headers = ["Feature", "Typical Projects", "Mart"]
    col_widths = [Inches(4.5), Inches(3.5), Inches(3.5)]
    col_x = [Inches(0.7)]
    for w in col_widths[:-1]:
        col_x.append(col_x[-1] + w + Inches(0.15))

    y_header = Inches(1.5)
    header_h = Inches(0.55)

    for i, (header, w, x) in enumerate(zip(headers, col_widths, col_x)):
        color = ACCENT_PRIMARY if i == 0 else (ACCENT_ROSE if i == 1 else ACCENT_EMERALD)
        hdr = add_rounded_rect(slide, x, y_header, w, header_h, color)
        set_text(hdr, header, font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Table rows
    comparison_rows = [
        ("Payment Integration", "❌  Mock / No payments", "✅  Razorpay + COD"),
        ("Order Tracking", "❌  No tracking", "✅  4-Stage live tracker"),
        ("Admin Dashboard", "❌  None or basic", "✅  Full CRUD + analytics"),
        ("UI / Design", "⚠️  Bootstrap / plain", "✅  Custom glassmorphic CSS"),
        ("Stock Management", "❌  Static data", "✅  Real-time auto-decrement"),
        ("Role-Based Auth", "⚠️  Single role only", "✅  User + Admin roles"),
        ("Profile Management", "❌  Not available", "✅  Photo upload + edit"),
        ("Database Setup", "⚠️  Manual SQL scripts", "✅  Auto-migrate on startup"),
        ("Deployment Ready", "❌  Localhost only", "✅  Render + Gunicorn"),
    ]

    row_h = Inches(0.52)
    start_y = y_header + header_h + Inches(0.08)

    for i, (feature, typical, mart) in enumerate(comparison_rows):
        y = start_y + i * (row_h + Inches(0.06))
        bg_color = BG_CARD if i % 2 == 0 else RGBColor(0x17, 0x20, 0x33)

        for j, (text, w, x) in enumerate(zip([feature, typical, mart], col_widths, col_x)):
            cell = add_rounded_rect(slide, x, y, w, row_h, bg_color)
            cell.line.fill.background()
            txt_color = WHITE if j == 0 else (RGBColor(0xF8, 0x71, 0x71) if "❌" in text or "⚠️" in text else ACCENT_EMERALD)
            add_text_box(slide, x + Inches(0.15), y + Inches(0.08), w - Inches(0.3), row_h - Inches(0.16),
                         text, font_size=12, color=txt_color, bold=(j == 0), alignment=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)

    add_gradient_bar(slide, Inches(0), Inches(7.46), SLIDE_WIDTH, Inches(0.04))


def build_code_highlight_slide(prs):
    """Slide 9 — Code Highlights"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_gradient_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.04))

    add_section_number(slide, Inches(0.8), Inches(0.4), "06", ACCENT_PRIMARY)
    add_text_box(slide, Inches(1.6), Inches(0.4), Inches(6), Inches(0.6),
                 "CODE HIGHLIGHTS", font_size=28, color=WHITE, bold=True)
    add_accent_bar(slide, Inches(1.6), Inches(1.05), Inches(1.5), Inches(0.04), ACCENT_PRIMARY)

    # Code block 1 — Payment Flow
    code1_title = "💳  Razorpay Payment Verification  (app.py)"
    code1 = """@app.route('/verify_payment', methods=['POST'])
def verify_payment():
    payment_id = request.form.get('razorpay_payment_id')
    order_id   = request.form.get('razorpay_order_id')
    signature  = request.form.get('razorpay_signature')
    
    params = {
        'razorpay_order_id':   order_id,
        'razorpay_payment_id': payment_id,
        'razorpay_signature':  signature
    }
    razorpay_client.utility.verify_payment_signature(params)
    # → Update order status, decrement stock, clear cart"""

    code_bg_1 = add_rounded_rect(slide, Inches(0.4), Inches(1.5), Inches(6.2), Inches(5.2), RGBColor(0x0D, 0x11, 0x1E), DARK_GRAY)
    add_text_box(slide, Inches(0.6), Inches(1.3), Inches(5.8), Inches(0.4),
                 code1_title, font_size=13, color=ACCENT_CYAN, bold=True)
    add_text_box(slide, Inches(0.65), Inches(1.85), Inches(5.7), Inches(4.5),
                 code1, font_size=10, color=RGBColor(0xA5, 0xF3, 0xFC), font_name='Consolas')

    # Code block 2 — Database Schema  
    code2_title = "🗄️  Auto-Migration Setup  (setup_db.py)"
    code2 = """def create_database_and_table(db_path='mart.db'):
    connection = sqlite3.connect(db_path)
    cursor = connection.cursor()

    cursor.execute(\"\"\"
    CREATE TABLE IF NOT EXISTS users (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        username VARCHAR(50) NOT NULL,
        email VARCHAR(100) NOT NULL UNIQUE,
        password_hash VARCHAR(255) NOT NULL,
        role VARCHAR(20) DEFAULT 'User',
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
    )\"\"\")
    
    # Auto-seed default data if tables empty
    cursor.execute("SELECT COUNT(*) FROM items")
    if cursor.fetchone()[0] == 0:
        cursor.executemany("INSERT INTO items ...", data)"""

    code_bg_2 = add_rounded_rect(slide, Inches(6.9), Inches(1.5), Inches(6.2), Inches(5.2), RGBColor(0x0D, 0x11, 0x1E), DARK_GRAY)
    add_text_box(slide, Inches(7.1), Inches(1.3), Inches(5.8), Inches(0.4),
                 code2_title, font_size=13, color=ACCENT_EMERALD, bold=True)
    add_text_box(slide, Inches(7.15), Inches(1.85), Inches(5.7), Inches(4.5),
                 code2, font_size=10, color=RGBColor(0xA5, 0xF3, 0xFC), font_name='Consolas')

    # Key Takeaways
    add_multi_text_box(slide, Inches(0.5), Inches(6.85), Inches(12), Inches(0.5), [
        {'text': '▸ Secure payment verification with cryptographic signatures     ▸ Zero-config database with auto-migration and seeding on first run',
         'size': 11, 'color': MID_GRAY, 'alignment': PP_ALIGN.LEFT},
    ])

    add_gradient_bar(slide, Inches(0), Inches(7.46), SLIDE_WIDTH, Inches(0.04))


def build_code_highlight_slide_2(prs):
    """Slide 10 — More Code Highlights (Cart & Admin)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_gradient_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.04))

    add_section_number(slide, Inches(0.8), Inches(0.4), "06", ACCENT_PRIMARY)
    add_text_box(slide, Inches(1.6), Inches(0.4), Inches(8), Inches(0.6),
                 "CODE HIGHLIGHTS — CART & ADMIN", font_size=28, color=WHITE, bold=True)
    add_accent_bar(slide, Inches(1.6), Inches(1.05), Inches(2), Inches(0.04), ACCENT_PRIMARY)

    # Code block — Smart Cart with Stock Validation
    code1_title = "🛒  Smart Cart — Stock-Aware Add to Cart  (app.py)"
    code1 = """@app.route('/add_to_cart/<int:item_id>', methods=['POST'])
def add_to_cart(item_id):
    quantity = int(request.form.get('quantity', 1))

    # Check stock availability
    cursor.execute("SELECT stock FROM items WHERE id=?", (item_id,))
    item = cursor.fetchone()
    if not item or item['stock'] < 1:
        flash("Item is out of stock!", "error")
        return redirect(url_for('dashboard'))

    # Update existing cart entry or insert new
    existing = cursor.execute(
        "SELECT * FROM cart WHERE user_id=? AND item_id=?",
        (session['user_id'], item_id)).fetchone()
    if existing:
        new_qty = min(existing['quantity'] + quantity, item['stock'])
        cursor.execute("UPDATE cart SET quantity=? WHERE id=?",
                       (new_qty, existing['id']))"""

    code_bg_1 = add_rounded_rect(slide, Inches(0.4), Inches(1.5), Inches(6.2), Inches(5.5), RGBColor(0x0D, 0x11, 0x1E), DARK_GRAY)
    add_text_box(slide, Inches(0.6), Inches(1.3), Inches(5.8), Inches(0.4),
                 code1_title, font_size=13, color=ACCENT_AMBER, bold=True)
    add_text_box(slide, Inches(0.65), Inches(1.85), Inches(5.7), Inches(5), 
                 code1, font_size=9.5, color=RGBColor(0xA5, 0xF3, 0xFC), font_name='Consolas')

    # Code block — Admin Dashboard Stats
    code2_title = "👨‍💼  Admin Dashboard Analytics  (app.py)"
    code2 = """@app.route('/admin_dashboard')
def admin_dashboard():
    if session.get('role') != 'Admin':
        flash("Access Denied. Admins only.", "error")
        return redirect(url_for('home'))

    # Real-time business metrics
    cursor.execute(
      "SELECT COUNT(*) as count FROM orders")
    total_orders = cursor.fetchone()['count']

    cursor.execute(
      "SELECT COALESCE(SUM(total_price), 0) as revenue "
      "FROM orders WHERE status NOT IN "
      "('pending','cancelled')")
    total_revenue = cursor.fetchone()['revenue']

    cursor.execute(
      "SELECT COUNT(*) as count "
      "FROM items WHERE stock < 5")
    low_stock = cursor.fetchone()['count']"""

    code_bg_2 = add_rounded_rect(slide, Inches(6.9), Inches(1.5), Inches(6.2), Inches(5.5), RGBColor(0x0D, 0x11, 0x1E), DARK_GRAY)
    add_text_box(slide, Inches(7.1), Inches(1.3), Inches(5.8), Inches(0.4),
                 code2_title, font_size=13, color=ACCENT_ROSE, bold=True)
    add_text_box(slide, Inches(7.15), Inches(1.85), Inches(5.7), Inches(5),
                 code2, font_size=9.5, color=RGBColor(0xA5, 0xF3, 0xFC), font_name='Consolas')

    add_gradient_bar(slide, Inches(0), Inches(7.46), SLIDE_WIDTH, Inches(0.04))


def build_thank_you_slide(prs):
    """Final Slide — Thank You"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_gradient_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06))

    # Decorative elements
    add_decorative_dots(slide, Inches(1), Inches(1), 6, 6, Inches(0.3), RGBColor(0x2A, 0x35, 0x50))
    add_decorative_dots(slide, Inches(10), Inches(5), 5, 5, Inches(0.3), RGBColor(0x2A, 0x35, 0x50))

    # Large circle decor
    big_circle = add_circle(slide, Inches(-2), Inches(4), Inches(6), RGBColor(0x1A, 0x22, 0x38))
    big_circle2 = add_circle(slide, Inches(10), Inches(-2), Inches(6), RGBColor(0x1A, 0x22, 0x38))

    # Thank you text
    add_text_box(slide, Inches(2), Inches(2.0), Inches(9.3), Inches(1),
                 "THANK YOU", font_size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Accent bar center
    add_accent_bar(slide, Inches(5.5), Inches(3.2), Inches(2.3), Inches(0.06), ACCENT_PRIMARY)

    # Subtitle
    add_multi_text_box(slide, Inches(2), Inches(3.6), Inches(9.3), Inches(1.5), [
        {'text': 'Mart — A Full-Stack E-Commerce Platform', 'size': 22, 'color': ACCENT_CYAN, 'bold': False, 'alignment': PP_ALIGN.CENTER, 'space_after': 16},
        {'text': 'Built with Flask • SQLite • Razorpay • Vanilla CSS', 'size': 14, 'color': MID_GRAY, 'alignment': PP_ALIGN.CENTER, 'space_after': 8},
        {'text': 'Questions?', 'size': 20, 'color': LIGHT_GRAY, 'bold': True, 'alignment': PP_ALIGN.CENTER},
    ])

    # Contact / info boxes
    info_items = [
        ("🔗", "GitHub Repository", "github.com/kushp9275", ACCENT_PRIMARY),
        ("🌐", "Live Demo", "e-commerce-website-demo.onrender.com", ACCENT_CYAN),
        ("📧", "Contact", "patelkartavya79@gmail.com", ACCENT_EMERALD),
    ]

    total_w = len(info_items) * Inches(3.5) + (len(info_items) - 1) * Inches(0.3)
    start_x = (Inches(13.333) - total_w) / 2

    for i, (emoji, label, value, color) in enumerate(info_items):
        x = start_x + i * (Inches(3.5) + Inches(0.3))
        y = Inches(5.3)
        card = add_rounded_rect(slide, x, y, Inches(3.5), Inches(1.0), BG_CARD, color)
        add_multi_text_box(slide, x + Inches(0.15), y + Inches(0.12), Inches(3.2), Inches(0.75), [
            {'text': f'{emoji}  {label}', 'size': 13, 'color': color, 'bold': True, 'alignment': PP_ALIGN.CENTER, 'space_after': 4},
            {'text': value, 'size': 10, 'color': LIGHT_GRAY, 'alignment': PP_ALIGN.CENTER},
        ])

    add_gradient_bar(slide, Inches(0), Inches(7.44), SLIDE_WIDTH, Inches(0.06))


# ─────────────────────────────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    print("Building slides...")
    build_title_slide(prs)        # 1
    build_agenda_slide(prs)       # 2
    build_overview_slide(prs)     # 3
    build_tech_stack_slide(prs)   # 4
    build_architecture_slide(prs) # 5
    build_er_diagram_slide(prs)   # 6
    build_usp_slide(prs)         # 7
    build_comparison_slide(prs)   # 8
    build_code_highlight_slide(prs)   # 9
    build_code_highlight_slide_2(prs) # 10
    build_thank_you_slide(prs)    # 11

    output_path = os.path.join(os.path.dirname(__file__), 'Mart_E-Commerce_Presentation.pptx')
    prs.save(output_path)
    print(f"[OK] Presentation saved to: {output_path}")
    print(f"   Total slides: {len(prs.slides)}")

if __name__ == '__main__':
    main()
